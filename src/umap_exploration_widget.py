
import os
import logging
from pathlib import Path
from typing import List, Optional, Tuple, Dict, Any
import datetime

import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
from matplotlib.widgets import LassoSelector, RectangleSelector, Button
from matplotlib.path import Path as MplPath
import ipywidgets as widgets
from IPython.display import display, clear_output

# Attempt to import scanpy and umap, handle if missing
try:
    import scanpy as sc
    import umap.umap_ as umap
    import scanpy as sc
    import umap.umap_ as umap
    from sklearn.cluster import KMeans
    from sklearn.preprocessing import StandardScaler
except ImportError:
    scanpy = None
    umap = None
    KMeans = None
    StandardScaler = None

# Import tile extraction helpers
# Assumes src/tile_extraction.py exists
try:
    from src.tile_extraction import (
        NotebookConfig, ChannelMapping, extract_multichannel_tile, 
        create_rgb_composite, export_tiles, create_default_mappings
    )
except ImportError:

    # Fallback for local development or if not yet in path
    try:
        from tile_extraction import (
            NotebookConfig, ChannelMapping, extract_multichannel_tile, 
            create_rgb_composite, export_tiles, create_default_mappings
        )
    except ImportError:
        logging.warning("Could not import tile_extraction. Tile generation features may be limited.")

# Import ChannelMappingWidget
try:
    from src.channel_mapping_widget import ChannelMappingWidget
except ImportError:
    try:
        from channel_mapping_widget import ChannelMappingWidget
    except ImportError:
        ChannelMappingWidget = None
        logging.warning("Could not import ChannelMappingWidget. Interactive tile config disabled.")


logger = logging.getLogger(__name__)

class UMAPExplorationWidget:
    """
    Widget for UMAP generation, clustering, exploration, and subsetting.
    """
    def __init__(self, df: pd.DataFrame, config: NotebookConfig, output_root: str = "UMAP"):
        self.df = df.copy().reset_index(drop=True)
        self.config = config
        self.output_root = Path(output_root)
        self.output_root.mkdir(parents=True, exist_ok=True)
        
        self.embeddings = None
        self.adata = None
        
        # State
        self.selected_indices = []
        self._points = None
        self._colors_default = None
        self.fig = None
        self.ax = None
        self.scatter = None
        self.lasso = None
        self.rect_selector = None
        self.current_subset_df = None
        
        self.current_subset_df = None
        
        # Init Figure once
        with plt.ioff(): # Prevent auto-display
             self.fig, self.ax = plt.subplots(figsize=(10, 7))
             self.fig.tight_layout(rect=[0, 0, 0.85, 1]) # Reserve space on right
        self.fig.canvas.header_visible = True
        self.fig.canvas.footer_visible = True
        
        # UI Elements
        self._create_widgets()
        
    def _create_widgets(self):
        # UMAP Parameters
        self.n_neighbors_slider = widgets.IntSlider(value=15, min=2, max=100, description='Neighbors:')
        self.min_dist_slider = widgets.FloatSlider(value=0.1, min=0.0, max=1.0, step=0.05, description='Min Dist:')
        self.metric_dropdown = widgets.Dropdown(options=['euclidean', 'cosine', 'manhattan'], value='euclidean', description='Metric:')
        self.n_neighbors_slider = widgets.IntSlider(value=15, min=2, max=100, description='Neighbors:')
        self.min_dist_slider = widgets.FloatSlider(value=0.1, min=0.0, max=1.0, step=0.05, description='Min Dist:')
        self.metric_dropdown = widgets.Dropdown(options=['euclidean', 'cosine', 'manhattan'], value='euclidean', description='Metric:')
        self.scale_data_checkbox = widgets.Checkbox(value=True, description='Scale Data', indent=False)
        self.random_state_input = widgets.IntText(value=42, description='Seed:', layout=widgets.Layout(width='150px'))
        self.run_umap_btn = widgets.Button(description='Generate UMAP', button_style='primary')
        self.run_umap_btn.on_click(self.run_umap)
        
        # Clustering Parameters
        self.leiden_res_slider = widgets.FloatSlider(value=0.5, min=0.1, max=2.0, step=0.1, description='Leiden Res:')
        self.run_leiden_btn = widgets.Button(description='Run Leiden')
        self.run_leiden_btn.on_click(self.run_leiden)
        
        self.kmeans_k_slider = widgets.IntSlider(value=5, min=2, max=20, description='KMeans K:')
        self.run_kmeans_btn = widgets.Button(description='Run KMeans')
        self.run_kmeans_btn.on_click(self.run_kmeans)
        
        # Visualization Controls
        self.color_mode_tgl = widgets.ToggleButtons(options=['Metadata', 'Feature'], value='Metadata', description='Color Mode:')
        self.color_mode_tgl.observe(self._on_color_mode_change, names='value')
        
        # Metadata Controls
        self.metadata_dropdown = widgets.Dropdown(description='Metadata:')
        self.metadata_dropdown.observe(self._on_color_change, names='value')
        
        # Feature Controls
        self.color_search_input = widgets.Text(description='Search Feature:', placeholder='Type to filter...')
        self.color_search_input.observe(self._on_color_search, names='value')
        
        self.feature_dropdown = widgets.Dropdown(description='Feature:')
        self.feature_dropdown.observe(self._on_color_change, names='value')
        
        self.log_scale_color_checkbox = widgets.Checkbox(value=False, description='Log Scale Color')
        self.log_scale_color_checkbox.observe(self._on_color_change, names='value')
        
        self.outlier_quantile_slider = widgets.FloatSlider(
            value=1.0, min=0.5, max=1.0, step=0.01, 
            description='Clip Max Q:', readout_format='.2f'
        )
        self.outlier_quantile_slider.observe(self._on_color_change, names='value')
        
        # --- NEW: Color Scheme Controls ---
        self.colormap_dropdown = widgets.Dropdown(
            options=['viridis', 'plasma', 'inferno', 'magma', 'cividis', 'coolwarm', 'RdBu_r', 'jet', 'Greys', 
                     'Blues', 'Reds', 'Greens', 'Purples', 'Oranges'],
            value='viridis',
            description='Colormap:'
        )
        self.colormap_dropdown.observe(self._on_color_change, names='value')
        
        self.palette_dropdown = widgets.Dropdown(
            options=['tab10', 'tab20', 'Set1', 'Dark2', 'Paired', 'Accent'],
            value='tab10',
            description='Palette:'
        )
        self.palette_dropdown.observe(self._on_color_change, names='value')
        
        self.clim_slider = widgets.FloatRangeSlider(
            value=[1.0, 99.0],
            min=0.0,
            max=100.0,
            step=0.5,
            description='Intensity %:',
            readout_format='.1f',
            layout=widgets.Layout(width='95%')
        )
        self.clim_slider.observe(self._on_color_change, names='value')
        
        self.point_size_slider = widgets.FloatSlider(
            value=5.0,
            min=0.1,
            max=50.0,
            step=0.5,
            description='Point Size:',
            readout_format='.1f'
        )
        self.point_size_slider.observe(self._on_color_change, names='value')
        
        # --- NEW: Filter/Visibility Controls ---
        self.filter_metadata_dropdown = widgets.Dropdown(description='Filter Meta:')
        self.filter_metadata_dropdown.observe(self._on_filter_meta_change, names='value')
        
        self.hide_values_select = widgets.SelectMultiple(description='Hide:', rows=3)
        self.hide_values_select.observe(self._on_visibility_change, names='value')
        
        self.fade_values_select = widgets.SelectMultiple(description='Fade Out:', rows=3)
        self.fade_values_select.observe(self._on_visibility_change, names='value')

        self.fade_alpha_slider = widgets.FloatSlider(value=0.1, min=0.0, max=1.0, step=0.05, description='Fade Alpha:')
        self.fade_alpha_slider.observe(self._on_visibility_change, names='value')
        # ----------------------------------
        

        
        self.subset_cluster_select = widgets.SelectMultiple(description='Subset Clust:', rows=5)
        
        # self.update_color_options() moved to end
        
        self.subset_cluster_btn = widgets.Button(description='Subset by Cluster')
        
        self.subset_cluster_btn.on_click(self.subset_by_cluster)
        
        # Selection & Actions
        self.selection_mode_tgl = widgets.ToggleButtons(options=['Lasso', 'Rectangle', 'None'], value='Lasso', description='Select Mode:')
        self.selection_mode_tgl.observe(self._on_selection_mode_change, names='value')
        
        self.save_name_input = widgets.Text(description='Save Name:', value=f"analysis_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}")
        self.save_format_dropdown = widgets.Dropdown(options=['CSV', 'Pickle (Faster)'], value='CSV', description='Save Format:')
        
        self.save_umap_btn = widgets.Button(description='Save UMAP Data')
        self.save_umap_btn.on_click(self.save_umap_data)
        
        self.load_umap_btn = widgets.Button(description='Load UMAP Data', button_style='info')
        self.load_umap_btn.on_click(self.load_umap_data)

        self.save_plot_name_input = widgets.Text(description='Plot Name:', value='umap_plot')
        self.save_plot_btn = widgets.Button(description='Save Plot')
        self.save_plot_btn.on_click(self.save_plot)
        
        self.export_format_dropdown = widgets.Dropdown(options=['Folder (Images)', 'PowerPoint (PPTX)'], value='PowerPoint (PPTX)', description='Format:')
        self.export_name_input = widgets.Text(description='Export Name:', value='exported_tiles', tooltip='Name for the output folder or PPTX file')
        self.max_tiles_input = widgets.IntText(value=50, description='Max Tiles:', layout=widgets.Layout(width='150px'))
        self.random_sample_checkbox = widgets.Checkbox(value=True, description='Random Sample', indent=False)
        
        self.generate_tiles_btn = widgets.Button(description='Generate Tiles (Selected)')
        self.generate_tiles_btn.on_click(self.generate_tiles_for_selection)
        
        self.create_subset_btn = widgets.Button(description='New UMAP from Selection')
        self.create_subset_btn.on_click(self.create_subset_from_selection)
        
        self.exclude_selection_btn = widgets.Button(description='New UMAP Exclude Selection')
        self.exclude_selection_btn.on_click(self.create_exclude_subset)
        
        self.drill_down_btn = widgets.Button(description='Drill Down (New Widget)', button_style='warning')
        self.drill_down_btn.on_click(self.drill_down)
        
        self.output_area = widgets.Output()
        self.tile_config_output = widgets.Output()
        self.plot_output = widgets.Output()
        self.drill_down_output = widgets.Output()
        
        # Initialize options last to ensure all widgets exist
        self.update_color_options()
        
        # Build Layout
        self._build_layout()

    def _build_layout(self):
        self.umap_box = widgets.VBox([
            widgets.HTML("<h3>UMAP Generation</h3>"),
            self.n_neighbors_slider, self.min_dist_slider, self.metric_dropdown, 
            widgets.HBox([self.scale_data_checkbox, self.random_state_input]), 
            self.run_umap_btn
        ])
        
        self.cluster_box = widgets.VBox([
            widgets.HTML("<h3>Clustering</h3>"),
            widgets.HBox([self.leiden_res_slider, self.run_leiden_btn]),
            widgets.HBox([self.kmeans_k_slider, self.run_kmeans_btn])
        ])
        
        self.metadata_controls = widgets.VBox([
            self.metadata_dropdown
        ])
        
        self.feature_controls = widgets.VBox([
            self.color_search_input,
            self.feature_dropdown,
            widgets.HBox([self.log_scale_color_checkbox, self.outlier_quantile_slider]),
            widgets.HBox([self.colormap_dropdown, self.clim_slider]) # Added clim and cmap
        ])
        
        # Initial visibility state
        self._update_controls_visibility()
        
        self.vis_box = widgets.VBox([
            widgets.HTML("<h3>Visualization & Selection</h3>"),
            self.color_mode_tgl,
            self.metadata_controls,
            self.feature_controls,
            widgets.HTML("<b>General Settings:</b>"),
            widgets.HBox([self.palette_dropdown, self.point_size_slider]), # Added palette and size
            widgets.HTML("<hr>"),
            widgets.HTML("<b>Data Visibility:</b>"),
            widgets.HBox([self.filter_metadata_dropdown, self.fade_alpha_slider]),
            widgets.HBox([self.hide_values_select, self.fade_values_select]),
            widgets.HTML("<hr>"),
            self.selection_mode_tgl,
            self.subset_cluster_select, self.subset_cluster_btn
        ])
        
        self.action_box = widgets.VBox([
            widgets.HTML("<h3>Actions & Export</h3>"),
            self.save_name_input,
            self.save_format_dropdown,
            widgets.HBox([self.save_umap_btn, self.load_umap_btn]),
            widgets.HBox([self.save_plot_name_input, self.save_plot_btn]),
            widgets.HTML("<hr>"),
            widgets.HTML("<b>Tile Export:</b>"),
            self.export_name_input,
            self.export_format_dropdown,
            widgets.HBox([self.max_tiles_input, self.random_sample_checkbox]),
            self.generate_tiles_btn,
            widgets.HTML("<hr>"),
            widgets.HBox([self.create_subset_btn, self.exclude_selection_btn]),
            self.drill_down_btn
        ])
        self.main_layout = widgets.HBox([
            widgets.VBox([self.umap_box, self.cluster_box, self.vis_box, self.action_box]),
            widgets.VBox([self.plot_output, self.output_area])
        ])
        
        self.app_layout = widgets.VBox([
            self.main_layout,
            self.drill_down_output,
            self.tile_config_output
        ])

    def display(self):
        display(self.app_layout)
        
        # Display the figure in its own output area or directly if using widget backend
        with self.plot_output:
             # Check if we are in an interactive backend
             import matplotlib
             backend = matplotlib.get_backend().lower()
             with self.output_area:
                 print(f"DEBUG: Matplotlib backend is '{backend}'")
                 
             if 'ipympl' in backend or 'widget' in backend:
                 if self.fig.canvas:
                     display(self.fig.canvas)
             else:
                 with self.output_area:
                     print("WARNING: Interactive mode not detected. Plot may be static.")
                     print("Ensure '%matplotlib widget' is run at the start of the notebook.")
                 # For static backends, we display the figure object itself
                 display(self.fig)

    def update_color_options(self):
        """Update options for metadata and feature dropdowns."""
        # 1. Metadata
        meta_options = [c for c in self.df.columns if c.startswith('Metadata_') or c in ['leiden', 'kmeans']]
        self.metadata_dropdown.options = ['None'] + sorted(meta_options)
        self.filter_metadata_dropdown.options = ['None'] + sorted(meta_options)
        
        # 2. Features
        feature_cols = self._get_feature_columns()
        all_features = sorted(list(set(feature_cols)))
        
        search_term = self.color_search_input.value.lower()
        if search_term:
            filtered = [o for o in all_features if search_term in o.lower()]
        else:
             filtered = all_features
             
        self.feature_dropdown.options = ['None'] + filtered
        
        # Preserve selection if possible
        if self.metadata_dropdown.value not in self.metadata_dropdown.options:
             self.metadata_dropdown.value = 'None'
             
        if self.feature_dropdown.value not in self.feature_dropdown.options:
             self.feature_dropdown.value = 'None'
             
        # Update cluster subset options
        cluster_cols = [c for c in self.df.columns if c in ['leiden', 'kmeans']]
        options = []
        for col in cluster_cols:
             vals = sorted(self.df[col].unique(), key=lambda x: int(x) if str(x).isdigit() else str(x))
             options.extend([f"{col}: {v}" for v in vals])
        self.subset_cluster_select.options = options

    def _get_feature_columns(self):
        """Exclude Metadata_ columns and coordinate columns."""
        exclude_cols = [c for c in self.df.columns if c.startswith('Metadata_')]
        exclude_cols += [
            self.config.x_column, self.config.y_column, 
            self.config.umap_x_column, self.config.umap_y_column,
            self.config.bbox_min_x_column, self.config.bbox_max_x_column,
            self.config.bbox_min_y_column, self.config.bbox_max_y_column,
            'ImageNumber', 'ObjectNumber'
        ]
        # Filter out None values
        exclude_cols = [c for c in exclude_cols if c is not None]
        feature_cols = [c for c in self.df.columns if c not in exclude_cols]
        # Ensure numeric
        numeric_cols = self.df[feature_cols].select_dtypes(include=[np.number]).columns.tolist()
        return numeric_cols

    def run_umap(self, btn):
        with self.output_area:
            print("Running UMAP... please wait.")
            
        try:
            feature_cols = self._get_feature_columns()
            if not feature_cols:
                with self.output_area:
                    print("Error: No feature columns found for UMAP.")
                return

            X = self.df[feature_cols].values
            
            if self.scale_data_checkbox.value and StandardScaler:
                with self.output_area: print(f"Scaling data ({X.shape[0]} cells, {X.shape[1]} features)...")
                scaler = StandardScaler()
                X = scaler.fit_transform(X)

            # Use scanpy if available for better integration, or direct umap
            import umap.umap_ as umap_module
            
            with self.output_area: print("Calculating UMAP embeddings (starting parallel execution)...")
            
            reducer = umap_module.UMAP(
                n_neighbors=self.n_neighbors_slider.value,
                min_dist=self.min_dist_slider.value,
                metric=self.metric_dropdown.value,
                random_state=self.random_state_input.value,
                n_jobs=-1
            )
            embedding = reducer.fit_transform(X)
            
            self.df['UMAP1'] = embedding[:, 0]
            self.df['UMAP2'] = embedding[:, 1]
            self.config.umap_x_column = 'UMAP1'
            self.config.umap_y_column = 'UMAP2'
            self.embeddings = embedding
            
            # Create AnnData for clustering
            if sc:
                 self.adata = sc.AnnData(X=X)
                 self.adata.obsm['X_umap'] = embedding
            
            with self.output_area:
                print("UMAP generation complete.")
            
            self.plot()
            
        except Exception as e:
            with self.output_area:
                print(f"Error running UMAP: {e}")

    def run_leiden(self, btn):
        if self.adata is None:
            with self.output_area: print("Please generate UMAP (and AnnData) first.")
            return
            
        with self.output_area: print("Running Leiden clustering...")
        try:
            # Requires neighbors graph first
            sc.pp.neighbors(self.adata, n_neighbors=self.n_neighbors_slider.value, use_rep='X')
            sc.tl.leiden(self.adata, resolution=self.leiden_res_slider.value, key_added='leiden')
            self.df['leiden'] = self.adata.obs['leiden'].values
            self.update_color_options()
            
            # Update subset dropdown
            self.update_color_options() # Logic moved to update_color_options func
            
            # Auto update plot if coloring by nothing
            if self.color_mode_tgl.value == 'Metadata' and self.metadata_dropdown.value == 'None':
                self.metadata_dropdown.value = 'leiden'
            else:
                self.plot() # Refresh plot to show new data if needed
                
            with self.output_area: print("Leiden clustering complete.")
        except Exception as e:
             with self.output_area: print(f"Error running Leiden: {e}")

    def run_kmeans(self, btn):
        with self.output_area: print("Running KMeans clustering...")
        try:
            # Run on embeddings or original features? Usually features for meaningful clusters, 
            # but sometimes embeddings for visual clusters. Let's do original features.
            feature_cols = self._get_feature_columns()
            X = self.df[feature_cols].values
            
            if self.scale_data_checkbox.value and StandardScaler:
                 scaler = StandardScaler()
                 X = scaler.fit_transform(X)
                 
            kmeans = KMeans(n_clusters=self.kmeans_k_slider.value, random_state=self.random_state_input.value).fit(X)
            self.df['kmeans'] = kmeans.labels_.astype(str)
            self.update_color_options()
            
            if self.color_mode_tgl.value == 'Metadata' and self.metadata_dropdown.value == 'None':
                self.metadata_dropdown.value = 'kmeans'
            else:
                self.plot()
                
            with self.output_area: print("KMeans complete.")
        except Exception as e:
             with self.output_area: print(f"Error running KMeans: {e}")

    def plot(self):
        if 'UMAP1' not in self.df.columns:
            return

        # Do NOT re-create figure. Clear Axes.
        self.ax.clear()
        
        # Determine visibility mask based on hidden values
        filter_col = self.filter_metadata_dropdown.value
        hide_vals = self.hide_values_select.value
        fade_vals = self.fade_values_select.value
        fade_alpha = self.fade_alpha_slider.value
        
        if filter_col and filter_col != 'None' and filter_col in self.df.columns and hide_vals:
             series_str = self.df[filter_col].astype(str)
             keep_mask = ~series_str.isin(hide_vals)
        else:
             keep_mask = np.ones(len(self.df), dtype=bool)

        # Store plotted indices to map back visual point indices to dataframe indices
        self.plotted_indices = np.where(keep_mask)[0]
        
        if len(self.plotted_indices) == 0:
            self.ax.set_title("No Data to Display (All points hidden)")
            self.fig.canvas.draw_idle()
            return
            
        plot_df = self.df.iloc[self.plotted_indices].copy()
        
        x = plot_df['UMAP1'].values
        y = plot_df['UMAP2'].values
        self._points = np.column_stack([x, y])
        
        color_col = 'None'
        mode = self.color_mode_tgl.value
        
        # with self.output_area: print(f"Plotting mode: {mode}")
        
        if mode == 'Metadata':
            color_col = self.metadata_dropdown.value
        else:
            color_col = self.feature_dropdown.value
            
        # with self.output_area: print(f"Update Plot: Mode={mode}, Col={color_col}")
        
        # Determine base alpha per point
        base_alphas = np.full(len(plot_df), 0.6)
        if filter_col and filter_col != 'None' and filter_col in plot_df.columns and fade_vals:
            series_str = plot_df[filter_col].astype(str)
            fade_mask = series_str.isin(fade_vals)
            base_alphas[fade_mask] = fade_alpha
            
        c_data = 'steelblue'
        
        if color_col != 'None' and color_col in plot_df.columns:
            c_series = plot_df[color_col].copy()
            point_size = self.point_size_slider.value
            
            # Numeric/Feature Mode
            if pd.api.types.is_numeric_dtype(c_series) and (mode == 'Feature' or c_series.nunique() > 20):
                # 1. Log Scale (Robust)
                label_suffix = ""
                if self.log_scale_color_checkbox.value:
                    min_val = c_series.min()
                    if min_val < 0:
                        c_series = c_series - min_val
                    c_series = np.log1p(c_series)
                    label_suffix = " (Log)"

                # 2. Robust Contrast (Percentile Clipping) using the new Range Slider
                # slider values are 0-100, we need 0.0-1.0 qt
                p_min, p_max = self.clim_slider.value
                
                v_min = c_series.quantile(p_min / 100.0) 
                v_max = c_series.quantile(p_max / 100.0)
                
                # Clip data
                c_data = c_series.clip(lower=v_min, upper=v_max)
                
                from mpl_toolkits.axes_grid1 import make_axes_locatable
                
                if not hasattr(self, 'cax'):
                    divider = make_axes_locatable(self.ax)
                    self.cax = divider.append_axes("right", size="5%", pad=0.05)
                
                self.cax.clear()
                self.cax.axis('on') # Ensure visible
                
                # Use selected colormap
                cmap_name = self.colormap_dropdown.value
                
                # Custom Linear Colormaps from Tab10
                # Tab10 indices: 0:Blue, 1:Orange, 2:Green, 3:Red, 4:Purple, 5:Brown, 6:Pink, 7:Gray, 8:Olive, 9:Cyan
                custom_cmaps = {
                    'Blues': ('#f7fbff', plt.get_cmap('tab10')(0)),   # Very light blue -> Blue
                    'Oranges': ('#fff5eb', plt.get_cmap('tab10')(1)), # Very light orange -> Orange
                    'Greens': ('#f7fcf5', plt.get_cmap('tab10')(2)),  # Very light green -> Green
                    'Reds': ('#fff5f0', plt.get_cmap('tab10')(3)),    # Very light red -> Red
                    'Purples': ('#fcfbfd', plt.get_cmap('tab10')(4))  # Very light purple -> Purple
                }
                
                if cmap_name in custom_cmaps:
                    from matplotlib.colors import LinearSegmentedColormap
                    c_start, c_end = custom_cmaps[cmap_name]
                    cmap = LinearSegmentedColormap.from_list(cmap_name, [c_start, c_end], N=256)
                else:
                    try:
                        cmap = plt.get_cmap(cmap_name)
                    except:
                        cmap = plt.get_cmap('viridis') # Fallback
                    
                from matplotlib.colors import Normalize
                from matplotlib.cm import ScalarMappable
                
                norm = Normalize(vmin=v_min, vmax=v_max)
                rgba_colors = cmap(norm(c_data))
                rgba_colors[:, 3] = base_alphas
                
                self.scatter = self.ax.scatter(x, y, c=rgba_colors, s=point_size)
                
                sm = ScalarMappable(norm=norm, cmap=cmap)
                sm.set_array([])
                self.fig.colorbar(sm, cax=self.cax, label=f"{color_col}{label_suffix}")
                self.cbar = True 
            
            # Categorical Mode
            else:
                cats = c_series.astype('category')
                if cats.isna().any():
                     if 'Unknown' not in cats.cat.categories:
                         cats = cats.cat.add_categories(['Unknown'])
                     cats = cats.fillna('Unknown')
                
                unique_cats = cats.cat.categories
                n_cats = len(unique_cats)
                
                # Use selected Palette
                palette_name = self.palette_dropdown.value
                try:
                    cmap = plt.get_cmap(palette_name)
                except:
                    cmap = plt.get_cmap('tab10')

                # Cycle colors
                # If n_cats > cmap.N (for qualitative colormaps like tab10 which have fixed N), 
                # we need to wrap around.
                # For continuous colormaps used as discrete, we might sample.
                # Assuming qualitative maps here.
                cat_to_color = {}
                for i, cat in enumerate(unique_cats):
                    if hasattr(cmap, 'colors'):
                         # ListedColormap
                         max_c = len(cmap.colors)
                         cat_to_color[cat] = cmap.colors[i % max_c]
                    else:
                         # LinearSegmentedColormap or similar
                         # Sample it? or just use tab20 fallback?
                         # Usually 'tab10', 'Set1' are Linear/Listed.
                         # Try accessing as callable.
                         cat_to_color[cat] = cmap(i % 20) # Modulo 20 just in case

                # Apply mapping
                from matplotlib.colors import to_rgba
                c_data_rgb = [to_rgba(cat_to_color[x]) for x in cats]
                c_data_rgba = np.array(c_data_rgb)
                c_data_rgba[:, 3] = base_alphas
                
                self.scatter = self.ax.scatter(x, y, c=c_data_rgba, s=point_size)
                
                # Legend
                handles = [
                    plt.Line2D([0], [0], marker='o', color='w', markerfacecolor=cat_to_color[cat], markersize=8, label=str(cat))
                    for cat in unique_cats
                ]
                # Pagination
                if len(handles) > 20: handles = handles[:20] 
                self.ax.legend(handles=handles, loc='center left', bbox_to_anchor=(1.02, 0.5), title=color_col)
                
                if hasattr(self, 'cax'):
                    self.cax.clear()
                    self.cax.axis('off')
        else:
            # No color selected
            point_size = self.point_size_slider.value
            from matplotlib.colors import to_rgba
            c_data_rgba = np.array([to_rgba('steelblue')] * len(plot_df))
            c_data_rgba[:, 3] = base_alphas
            self.scatter = self.ax.scatter(x, y, c=c_data_rgba, s=point_size)
            if hasattr(self, 'cax'):
                self.cax.clear()
                self.cax.axis('off')

        self.ax.set_xlabel('UMAP1')
        self.ax.set_ylabel('UMAP2')
        self.ax.set_title(f'UMAP Projection ({len(self.df)} cells)')
        
        self._colors_default = self.scatter.get_facecolors().copy()
        if len(self._colors_default) == 1 and len(self._points) > 1:
            self._colors_default = np.repeat(self._colors_default, len(self._points), axis=0)

        # Init selectors on the specific AX
        self._init_selectors()
        
        # Refresh canvas
        self.fig.canvas.draw_idle()
        
        # If not using an interactive backend, we must re-display the figure/canvas explicitly
        import matplotlib
        backend = matplotlib.get_backend().lower()
        if 'ipympl' not in backend and 'widget' not in backend:
             with self.plot_output:
                 clear_output(wait=True)
                 display(self.fig)

    def _init_selectors(self):
        # Cleanup old selectors if they exist
        if self.lasso:
            self.lasso.set_active(False)
            self.lasso = None
        if self.rect_selector:
            self.rect_selector.set_active(False)
            self.rect_selector = None
            
        # Re-initialize with useblit=False for better stability in notebooks
        self.lasso = LassoSelector(self.ax, self._on_select, useblit=False, button=1)
        self.rect_selector = RectangleSelector(
            self.ax, self._on_rect_select, useblit=False, button=1,
            minspanx=5, minspany=5, spancoords='pixels', interactive=True
        )
        self._set_selector_active()

    def _set_selector_active(self):
        # Force disable Pan/Zoom tools if they are active in the toolbar
        try:
            # For ipympl/widget backend
            if self.fig.canvas.toolbar:
                # This depends on the backend implementation, but usually setting mode to '' works
                # or calling the pan/zoom methods to toggle them off if active.
                # Taking a safer approach: just ensure we try to reset.
                if hasattr(self.fig.canvas.toolbar, 'mode') and self.fig.canvas.toolbar.mode != '':
                     # print("Disabling toolbar mode for selection...")
                     self.fig.canvas.toolbar.mode = '' 
        except Exception:
            pass
            
        mode = self.selection_mode_tgl.value
        if mode == 'Lasso':
            self.lasso.set_active(True)
            self.rect_selector.set_active(False)
        elif mode == 'Rectangle':
            self.lasso.set_active(False)
            self.rect_selector.set_active(True)
        else:
            self.lasso.set_active(False)
            self.rect_selector.set_active(False)

    def _on_selection_mode_change(self, change):
        if self.lasso:
            self._set_selector_active()

    def _on_select(self, vertices):
        path = MplPath(vertices)
        mask = path.contains_points(self._points)
        selected_plot_indices = np.where(mask)[0]
        if hasattr(self, 'plotted_indices'):
            self.selected_indices = self.plotted_indices[selected_plot_indices].tolist()
        else:
            self.selected_indices = selected_plot_indices.tolist()
        self._highlight_selection(selected_plot_indices)

    def _on_rect_select(self, eclick, erelease):
        x1, y1 = eclick.xdata, eclick.ydata
        x2, y2 = erelease.xdata, erelease.ydata
        xmin, xmax = min(x1, x2), max(x1, x2)
        ymin, ymax = min(y1, y2), max(y1, y2)
        
        mask = (
            (self._points[:, 0] >= xmin) & (self._points[:, 0] <= xmax) &
            (self._points[:, 1] >= ymin) & (self._points[:, 1] <= ymax)
        )
        selected_plot_indices = np.where(mask)[0]
        if hasattr(self, 'plotted_indices'):
            self.selected_indices = self.plotted_indices[selected_plot_indices].tolist()
        else:
            self.selected_indices = selected_plot_indices.tolist()
        self._highlight_selection(selected_plot_indices)

    def _highlight_selection(self, plot_indices=None):
        if self.scatter is None: return
        colors = self._colors_default.copy()
        
        if plot_indices is None:
            if hasattr(self, 'plotted_indices'):
                mask = np.isin(self.plotted_indices, self.selected_indices)
                plot_indices = np.where(mask)[0]
            else:
                plot_indices = self.selected_indices
                
        if len(plot_indices) > 0:
            colors[plot_indices] = [1, 0, 0, 1] # Red highlight
            with self.output_area:
                print(f"Selected {len(self.selected_indices)} points.")
        
        self.scatter.set_facecolors(colors)
        self.fig.canvas.draw_idle()

    def _on_color_search(self, change):
        try:
             self.update_color_options()
        except Exception as e:
            with self.output_area:
                print(f"Error in color search: {e}")

    def _on_color_mode_change(self, change):
        try:
            with self.output_area:
                # print(f"Debug: Color mode changed to {change['new']}")
                self._update_controls_visibility()
                self.plot()
        except Exception as e:
            with self.output_area:
                print(f"Error in color mode change: {e}")
                import traceback
                traceback.print_exc()

    def _update_controls_visibility(self):
        mode = self.color_mode_tgl.value
        if mode == 'Metadata':
            self.metadata_controls.layout.display = 'block'
            self.feature_controls.layout.display = 'none'
        else:
            self.metadata_controls.layout.display = 'none'
            self.feature_controls.layout.display = 'block'

    def _on_color_change(self, change):
        try:
            self.plot()
        except Exception as e:
            with self.output_area:
                print(f"Error in color change: {e}")
                import traceback
                traceback.print_exc()
                
    def _on_filter_meta_change(self, change):
        try:
            meta_col = self.filter_metadata_dropdown.value
            if meta_col and meta_col != 'None' and meta_col in self.df.columns:
                unique_vals = [str(x) for x in self.df[meta_col].unique() if pd.notna(x)]
                sorted_vals = sorted(unique_vals)
                self.hide_values_select.options = sorted_vals
                self.fade_values_select.options = sorted_vals
            else:
                self.hide_values_select.options = []
                self.fade_values_select.options = []
            
            self.hide_values_select.value = ()
            self.fade_values_select.value = ()
            self.plot()
        except Exception as e:
            with self.output_area:
                print(f"Error in filter meta change: {e}")

    def _on_visibility_change(self, change):
        try:
            self.plot()
        except Exception as e:
            with self.output_area:
                print(f"Error in visibility change: {e}")

    def save_umap_data(self, btn):
        name = self.save_name_input.value
        if not name:
             with self.output_area: print("Please enter a save name.")
             return
             
        save_dir = self.output_root / name
        save_dir.mkdir(parents=True, exist_ok=True)
        
        format_choice = self.save_format_dropdown.value
        if format_choice == 'CSV':
            csv_path = save_dir / "umap_data.csv"
            self.df.to_csv(csv_path, index=False)
            with self.output_area:
                print(f"Saved UMAP data to {csv_path}")
        else:
            pkl_path = save_dir / "umap_data.pkl"
            self.df.to_pickle(pkl_path)
            with self.output_area:
                print(f"Saved UMAP data to {pkl_path}")

    def load_umap_data(self, btn):
        """Load UMAP data from a CSV or Pickle file based on the save name."""
        name = self.save_name_input.value
        if not name:
             with self.output_area: print("Please enter the name of the folder to load from in the 'Save Name' box.")
             return
             
        load_dir = self.output_root / name
        if not load_dir.exists():
            with self.output_area: print(f"Directory {load_dir} does not exist.")
            return
            
        csv_path = load_dir / "umap_data.csv"
        pkl_path = load_dir / "umap_data.pkl"
        
        try:
            if pkl_path.exists():
                with self.output_area: print(f"Loading {pkl_path}...")
                self.df = pd.read_pickle(pkl_path)
            elif csv_path.exists():
                with self.output_area: print(f"Loading {csv_path}...")
                self.df = pd.read_csv(csv_path)
            else:
                with self.output_area: print(f"No umap_data.pkl or umap_data.csv found in {load_dir}.")
                return
                
            # If UMAP embeddings exist, ensure config points to them
            if 'UMAP1' in self.df.columns and 'UMAP2' in self.df.columns:
                self.config.umap_x_column = 'UMAP1'
                self.config.umap_y_column = 'UMAP2'
                self.embeddings = self.df[['UMAP1', 'UMAP2']].values
            else:
                 with self.output_area: print("Warning: Loaded data does not contain UMAP1/UMAP2 columns. You may need to run UMAP.")
            
            # Reconstruct AnnData for clustering if scanpy is available
            try:
                import scanpy as sc
                feature_cols = self._get_feature_columns()
                if feature_cols:
                    X = self.df[feature_cols].values
                    self.adata = sc.AnnData(X=X)
                    if self.embeddings is not None:
                        self.adata.obsm['X_umap'] = self.embeddings
            except Exception as e:
                pass # scanpy not available or no features, ignore
                 
            self.update_color_options()
            self.plot()
            with self.output_area: print("Load complete!")
        except Exception as e:
            with self.output_area: print(f"Error loading UMAP data: {e}")

    def save_plot(self, btn):
        if self.fig is None:
            with self.output_area: print("No plot to save.")
            return
            
        name = self.save_name_input.value
        plot_name = self.save_plot_name_input.value
        
        if not name:
             with self.output_area: print("Please enter a save folder name.")
             return
        if not plot_name:
             with self.output_area: print("Please enter a plot file name.")
             return

        save_dir = self.output_root / name
        save_dir.mkdir(parents=True, exist_ok=True)
        
        if not plot_name.endswith('.png'):
            plot_name += '.png'
            
        plot_path = save_dir / plot_name
        self.fig.savefig(plot_path, dpi=300, bbox_inches='tight')
         
        with self.output_area:
            print(f"Saved plot to {plot_path}")


    def generate_tiles_for_selection(self, btn):
        if not self.selected_indices:
            with self.output_area: print("No selection to generate tiles for.")
            return

        # Apply Random Sampling if requested
        export_indices = self.selected_indices
        
        if self.random_sample_checkbox.value:
            max_tiles = self.max_tiles_input.value
            if len(self.selected_indices) > max_tiles:
                import random
                export_indices = random.sample(self.selected_indices, max_tiles)
                with self.output_area:
                    print(f"Randomly sampled {max_tiles} tiles from {len(self.selected_indices)} selected.")
            else:
                 with self.output_area:
                    print(f"Selection size ({len(self.selected_indices)}) is within max tiles limit. Exporting all.")
        
        # Store for export
        self._export_indices = export_indices

        with self.tile_config_output:
            clear_output()
            if ChannelMappingWidget is None:
                 print("ChannelMappingWidget not available. Cannot configure.")
                 return
                 
            print("Extracting sample tiles for configuration...")
            # Extract samples (up to 5 random from the *export* set)
            import random
            n_samples = min(5, len(export_indices))
            sample_indices_preview = random.sample(export_indices, n_samples)
            
            sample_tiles = []
            try:
                for idx in sample_indices_preview:
                    row = self.df.iloc[idx]
                    tile = extract_multichannel_tile(row, self.config)
                    sample_tiles.append(tile)
                
                # Init Widget
                self.cm_widget = ChannelMappingWidget(self.config, sample_tiles)
                
                confirm_btn = widgets.Button(description="Confirm & Export Tiles", button_style='success')
                confirm_btn.on_click(self._on_confirm_export)
                
                display(widgets.VBox([
                    widgets.HTML("<h4>Configure Tile Export</h4>"),
                    self.cm_widget.widget_container,
                    confirm_btn
                ]))
                
            except Exception as e:
                print(f"Error extracting sample or initializing widget: {e}")
                import traceback
                traceback.print_exc()

    def _on_confirm_export(self, btn):
        name = self.save_name_input.value
        export_name = self.export_name_input.value
        
        if not export_name:
             with self.output_area: print("Please enter an export name.")
             return
             
        # Use the stored indices from the generation step
        if hasattr(self, '_export_indices'):
            indices = self._export_indices
        else:
            indices = self.selected_indices
            
        subset = self.df.iloc[indices]
        mappings = self.cm_widget.get_mappings()
        
        format_val = self.export_format_dropdown.value
        
        if 'PowerPoint' in format_val:
            # PPTX Export
            try:
                if not export_name.endswith('.pptx'):
                    export_name += '.pptx'
                pptx_path = self.output_root / name / export_name
                pptx_path.parent.mkdir(parents=True, exist_ok=True)
                
                with self.output_area: print(f"Exporting {len(subset)} tiles to PowerPoint... ({pptx_path})")
                
                self._export_to_pptx(subset, self.config, mappings, pptx_path)
                
                with self.output_area: print(f"PPTX Export complete: {pptx_path}")
                with self.tile_config_output:
                     clear_output()
                     print("Export complete.")

            except Exception as e:
                with self.output_area:
                    print(f"PPTX Export failed: {e}")
                    import traceback
                    traceback.print_exc()
        else:
            # Standard Folder Export
            tiles_dir = self.output_root / name / export_name
            tiles_dir.mkdir(parents=True, exist_ok=True)
            
            with self.output_area: print("Generating tiles... check console/logs for progress.")
            
            # Update config output dir
            export_config = NotebookConfig(
                csv_path=self.config.csv_path,
                image_base_path=self.config.image_base_path,
                output_dir=str(tiles_dir),
                tile_size=self.config.tile_size,
                x_column=self.config.x_column,
                y_column=self.config.y_column,
                umap_x_column=self.config.umap_x_column,
                umap_y_column=self.config.umap_y_column,
                well_column=self.config.well_column,
                field_column=self.config.field_column,
                channel_names=self.config.channel_names,
                filename_pattern=self.config.filename_pattern,
                bbox_min_x_column=self.config.bbox_min_x_column,
                bbox_min_y_column=self.config.bbox_min_y_column,
                bbox_max_x_column=self.config.bbox_max_x_column,
                bbox_max_y_column=self.config.bbox_max_y_column
            )
            
            try:
                res = export_tiles(subset, export_config, mappings, output_format='png')
                with self.output_area:
                    print(f"Exported {res['successful']} tiles to {res['output_path']}")
                    if res['failed'] > 0:
                        print(f"Failed to export {res['failed']} tiles.")
                        print("First 5 errors:")
                        for err in res['errors'][:5]:
                            print(f"  - {err}")
                
                # Clear config area
                with self.tile_config_output:
                    clear_output()
                    print("Export complete.")
            except Exception as e:
                with self.output_area:
                    print(f"Export failed: {e}")

    def _export_to_pptx(self, subset_df, config, mappings, output_path):
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            raise ImportError("python-pptx is not installed. Please run: pip install python-pptx")
            
        prs = Presentation()
        
        # Define layout (Blank)
        # Layouts: 0=Title, 1=Title+Content, 6=Blank
        try:
            blank_slide_layout = prs.slide_layouts[6] 
        except:
             blank_slide_layout = prs.slide_layouts[0] # Fallback
        
        from io import BytesIO
        from PIL import Image as PILImage
        
        count = 0
        total = len(subset_df)
        
        for idx, row in subset_df.iterrows():
            try:
                # 1. Generate Tile (Numpy Array)
                tile = extract_multichannel_tile(row, config)
                
                # 2. Convert to RGB Composite
                rgb_tile = create_rgb_composite(tile, mappings)
                
                # 3. Convert to PIL Image
                pil_img = PILImage.fromarray(rgb_tile)
                
                # 4. Save to BytesIO for PPTX
                img_stream = BytesIO()
                pil_img.save(img_stream, format='PNG')
                img_stream.seek(0)
                
                # 5. Add Slide
                slide = prs.slides.add_slide(blank_slide_layout)
                
                # 6. Add Image (centered, reasonably large)
                # Slide is usually 10x7.5 inches
                left = Inches(1)
                top = Inches(1.5)
                height = Inches(5)
                pic = slide.shapes.add_picture(img_stream, left, top, height=height)
                
                # 7. Add Title / Metadata
                # Get some useful ID
                well = str(row.get(config.well_column, 'Unknown'))
                field = str(row.get(config.field_column, 'Unknown'))
                obj_id = str(row.get('ObjectNumber', 'Unknown'))
                
                title_text = f"Well: {well}, Field: {field}, Obj: {obj_id}"
                
                # Add text box
                title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.5))
                tf = title_box.text_frame
                tf.text = title_text
                
                count += 1
                if count % 10 == 0:
                     with self.output_area:
                         print(f"Processed {count}/{total} slides...", end='\r')
                         
            except Exception as e:
                print(f"Skipping tile for row {idx}: {e}")
                
        prs.save(output_path)

    def create_subset_from_selection(self, btn):
        if not self.selected_indices: return
        self._create_subset_umap(self.selected_indices, f"{self.save_name_input.value}_subset")

    def create_exclude_subset(self, btn):
        if not self.selected_indices: return
        all_indices = set(range(len(self.df)))
        keep_indices = list(all_indices - set(self.selected_indices))
        self._create_subset_umap(keep_indices, f"{self.save_name_input.value}_excluded")
        
    def subset_by_cluster(self, btn):
        selections = self.subset_cluster_select.value
        if not selections: return
        
        # Handle multiple selections
        # Format "method: cluster_id"
        
        combined_indices = []
        name_parts = []
        
        for val in selections:
            method, cluster_id = val.split(': ')
            mask = self.df[method].astype(str) == cluster_id
            indices = np.where(mask)[0].tolist()
            combined_indices.extend(indices)
            if cluster_id not in name_parts:
                name_parts.append(cluster_id)
        
        combined_indices = sorted(list(set(combined_indices)))
        
        suffix = "_".join(name_parts[:3])
        if len(name_parts) > 3: suffix += "_etc"
        
        self._create_subset_umap(combined_indices, f"{self.save_name_input.value}_clusters_{suffix}")

    def _create_subset_umap(self, indices, name):
        subset_df = self.df.iloc[indices].copy()
        save_dir = self.output_root / name
        save_dir.mkdir(parents=True, exist_ok=True)
        
        # Save the subset data so it can be loaded later
        csv_path = save_dir / "umap_data.csv"
        subset_df.to_csv(csv_path, index=False)
        
        with self.output_area:
             print(f"Created subset '{name}' with {len(subset_df)} cells.")
             print("To analyze this subset, verify it is saved.")

    def drill_down(self, btn):
        """Spawn a new UMAP widget for the selected subset."""
        if not self.selected_indices:
             with self.output_area: print("No selection to drill down into.")
             return
             
        n_selected = len(self.selected_indices)
        with self.output_area: 
            print(f"Drilling down into {n_selected} cells... check below.")
        
        # 1. Create Subset
        subset_df = self.df.iloc[self.selected_indices].reset_index(drop=True)
        
        # 2. Spawn Child Widget
        # We use the same config and output root
        self.child_widget = UMAPExplorationWidget(subset_df, self.config, output_root=self.output_root)
        
        # Pre-set default save name
        self.child_widget.save_name_input.value = f"{self.save_name_input.value}_subset"
        
        # 3. Display in Drill Down Area
        with self.drill_down_output:
            clear_output()
            display(widgets.HTML(f"<h4>Drill Down: {n_selected} cells</h4>"))
            self.child_widget.display()


class UMAPViewerWidget:
    """
    Widget to view saved UMAP analyses.
    """
    def __init__(self, base_dir: str = "UMAP"):
        self.base_dir = Path(base_dir)
        self.df = None
        
        # State
        self.fig = None
        self.ax = None
        self.scatter = None
        self._points = None
        self._colors_default = None

        # Init Figure once
        with plt.ioff(): 
             self.fig, self.ax = plt.subplots(figsize=(8, 8))
        self.fig.canvas.header_visible = False
        self.fig.canvas.footer_visible = False
        
        self._create_widgets()
        
    def _create_widgets(self):
        self.session_dropdown = widgets.Dropdown(description='Session:')
        self.load_btn = widgets.Button(description='Load')
        self.load_btn.on_click(self.load_session)
        
        self.color_by_dropdown = widgets.Dropdown(description='Color By:')
        self.color_by_dropdown.observe(self._on_color_change, names='value')
        
        self.output_area = widgets.Output() # For logs
        self.plot_output = widgets.Output() # For plot
        
        self.refresh_sessions()
        
    def refresh_sessions(self):
        if self.base_dir.exists():
            sessions = [d.name for d in self.base_dir.iterdir() if d.is_dir()]
            self.session_dropdown.options = sorted(sessions)
        else:
            self.session_dropdown.options = []

    def update_color_options(self):
        if self.df is None: return
        options = [c for c in self.df.columns if c.startswith('Metadata_') or c in ['leiden', 'kmeans']]
        self.color_by_dropdown.options = ['None'] + sorted(options)
        
        # Auto-select a cluster column if available and currently None
        if self.color_by_dropdown.value == 'None' or not self.color_by_dropdown.value:
            if 'leiden' in self.df.columns:
                self.color_by_dropdown.value = 'leiden'
            elif 'kmeans' in self.df.columns:
                self.color_by_dropdown.value = 'kmeans'

    def _on_color_change(self, change):
        if self.df is not None:
            self.plot()

    def load_session(self, btn):
        session = self.session_dropdown.value
        if not session: return
        
        path = self.base_dir / session / "umap_data.csv"
        if not path.exists():
            with self.output_area: print(f"No data found at {path}")
            return
            
        with self.output_area:
            print(f"Loading {session}...")
            self.df = pd.read_csv(path)
            
            # Ensure UMAP cols exist or fallback
            if 'UMAP1' not in self.df.columns:
                 # Try last two cols? Or just error
                 cols = self.df.select_dtypes(include=[np.number]).columns
                 if len(cols) >= 2:
                     self.df['UMAP1'] = self.df[cols[-2]]
                     self.df['UMAP2'] = self.df[cols[-1]]
            
            self.update_color_options()
            self.plot()
            print(f"Loaded {len(self.df)} cells.")

    def plot(self):
        if self.df is None: return
        
        x = self.df['UMAP1'].values
        y = self.df['UMAP2'].values
        
        self.ax.clear()
        
        color_col = self.color_by_dropdown.value
        c_data = 'steelblue'
        
        if color_col and color_col != 'None' and color_col in self.df.columns:
            c_series = self.df[color_col]
            if pd.api.types.is_numeric_dtype(c_series) and c_series.nunique() > 20:
                    c_data = c_series
                    cmap = 'viridis'
                    self.scatter = self.ax.scatter(x, y, c=c_data, s=5, cmap=cmap, alpha=0.6)
                    plt.colorbar(self.scatter, ax=self.ax, label=color_col)
            else:
                # Categorical
                cats = c_series.astype('category')
                
                # Handle NaNs
                if cats.isna().any():
                     if 'Unknown' not in cats.cat.categories:
                         cats = cats.cat.add_categories(['Unknown'])
                     cats = cats.fillna('Unknown')
                
                unique_cats = cats.cat.categories
                cmap = plt.get_cmap('tab20')
                cat_to_color = {cat: cmap(i % 20) for i, cat in enumerate(unique_cats)}
                
                # List comprehension mapping (Safe)
                c_data = [cat_to_color[x] for x in cats]
                
                self.scatter = self.ax.scatter(x, y, c=c_data, s=5, alpha=0.6)
                
                # Legend
                handles = [
                    plt.Line2D([0], [0], marker='o', color='w', markerfacecolor=cat_to_color[cat], markersize=8, label=cat)
                    for cat in unique_cats
                ]
                if len(handles) > 20: handles = handles[:20] 
                self.ax.legend(handles=handles, loc='center left', bbox_to_anchor=(1, 0.5), title=color_col)
        else:
             self.scatter = self.ax.scatter(x, y, c=c_data, s=5, alpha=0.6)
             
        self.ax.set_xlabel('UMAP1')
        self.ax.set_ylabel('UMAP2')
        self.ax.set_title(f"View: {self.session_dropdown.value} ({len(self.df)} cells)")
        
        self.fig.canvas.draw_idle()
        
        # Redraw for static backends
        import matplotlib
        backend = matplotlib.get_backend().lower()
        if 'ipympl' not in backend and 'widget' not in backend:
             with self.plot_output:
                 clear_output(wait=True)
                 display(self.fig)

    def display(self):
        controls = widgets.HBox([
            self.session_dropdown, 
            self.load_btn, 
            self.color_by_dropdown
        ])
        
        display(widgets.VBox([controls, self.output_area]))
        
        with self.plot_output:
             import matplotlib
             backend = matplotlib.get_backend().lower()
             if 'ipympl' in backend or 'widget' in backend:
                 display(self.fig.canvas)
        display(self.plot_output)
